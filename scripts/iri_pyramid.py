from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Helper to convert pixel-like coordinates to pptx EMUs
def px(val):
    return int(val * 9525)  # PowerPoint uses EMUs; 1 px ≈ 9525 EMUs

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

# --- Shape definitions from your spec ---
shapes_spec = [
    {
        "type": "trapezoid",
        "x": 200, "y": 650, "w": 1520, "h": 180,
        "color": RGBColor(0x4C, 0xAF, 0x50),  # Green
        "opacity": 0.85,
        "text": "Measurement & Validation:\nIRI foundations, indicators, weighting, validation, pilot dataset"
    },
    {
        "type": "trapezoid",
        "x": 260, "y": 520, "w": 1400, "h": 150,
        "color": RGBColor(0xFF, 0x98, 0x00),  # Orange
        "opacity": 0.85,
        "text": "Comparative & Spatial Analysis:\nMulti-country IRI, sub-national IRI, QGIS mapping"
    },
    {
        "type": "trapezoid",
        "x": 320, "y": 400, "w": 1280, "h": 130,
        "color": RGBColor(0x67, 0x3A, 0xB7),  # Purple
        "opacity": 0.85,
        "text": "Predictive & Scenario Modelling:\nOutage prediction, Monte Carlo, stress-testing"
    },
    {
        "type": "trapezoid",
        "x": 380, "y": 290, "w": 1160, "h": 110,
        "color": RGBColor(0x00, 0x96, 0x88),  # Teal
        "opacity": 0.85,
        "text": "Sectoral & Sovereignty Extensions:\nE-IRI, Cloud/AI, vendor dependence"
    },
    {
        "type": "triangle",
        "x": 760, "y": 120, "w": 400, "h": 180,
        "color": RGBColor(0x0D, 0x47, 0xA1),  # Deep Blue
        "opacity": 0.90,
        "text": "Synthesis & Institutionalization:\nIRI 2.0, longitudinal analysis, Atlas, dashboard, monograph"
    }
]

# --- Draw shapes ---
for spec in shapes_spec:
    if spec["type"] == "trapezoid":
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE (we will convert to trapezoid-like)
            px(spec["x"]), px(spec["y"]),
            px(spec["w"]), px(spec["h"])
        )
    else:
        shape = slide.shapes.add_shape(
            3,  # MSO_SHAPE.ISOSCELES_TRIANGLE
            px(spec["x"]), px(spec["y"]),
            px(spec["w"]), px(spec["h"])
        )

    # Fill color
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = spec["color"]
    fill.transparency = 1 - spec["opacity"]

    # Text
    tf = shape.text_frame
    tf.text = spec["text"]
    for p in tf.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = "Segoe UI"
        p.alignment = 1  # center

# --- Footer ---
footer = slide.shapes.add_textbox(px(400), px(860), px(1200), px(60))
tf = footer.text_frame
tf.text = "A scalable, multi-year research program for digital resilience and sovereignty."
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(85, 85, 85)
p.font.name = "Segoe UI"

prs.save("iri_pyramid_generated.pptx")
