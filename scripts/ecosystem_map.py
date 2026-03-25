from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Helper to convert pixel-like coordinates to pptx EMUs
def px(val):
    return int(val * 9525)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

# --- Central Node ---
central = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    px(760), px(300), px(400), px(200)
)
fill = central.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x0D, 0x47, 0xA1)  # Deep blue
tf = central.text_frame
tf.text = "Digital Infrastructure Ecosystem\n(Core networks, IXPs, cables, data centers, cloud regions)"
for p in tf.paragraphs:
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

# --- Stakeholder Nodes ---
nodes = [
    {
        "label": "Regulators & Authorities\n- ICT regulators\n- Competition bodies\n- Licensing agencies\n- Data protection authorities",
        "x": 300, "y": 80,
        "color": RGBColor(0x4C, 0xAF, 0x50)  # Green
    },
    {
        "label": "Ministries & Government\n- ICT/Digital Economy\n- Planning commissions\n- Security agencies\n- Emergency response units",
        "x": 1200, "y": 80,
        "color": RGBColor(0x00, 0x96, 0x88)  # Teal
    },
    {
        "label": "IXPs & Operators\n- IXPs\n- Mobile operators\n- Fiber providers\n- Cable consortia",
        "x": 1450, "y": 350,
        "color": RGBColor(0xFF, 0x98, 0x00)  # Orange
    },
    {
        "label": "Cloud & Data Providers\n- Hyperscalers\n- CDN providers\n- Data centers\n- Platforms",
        "x": 1200, "y": 650,
        "color": RGBColor(0x67, 0x3A, 0xB7)  # Purple
    },
    {
        "label": "Election & Civic Institutions\n- Election commissions\n- Voter registration\n- Transmission systems\n- Civic tech groups",
        "x": 300, "y": 650,
        "color": RGBColor(0xFF, 0xC1, 0x07)  # Gold
    }
]

for node in nodes:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        px(node["x"]), px(node["y"]),
        px(350), px(200)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = node["color"]

    tf = shape.text_frame
    tf.text = node["label"]
    for p in tf.paragraphs:
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER

# --- Connector Lines (straight lines for clarity) ---
connections = [
    (300+175, 80+100, 760, 300+100),      # Regulators → Central
    (1200+175, 80+100, 760+400, 300+100), # Ministries → Central
    (1450, 350+100, 760+400, 300+100),    # IXPs → Central
    (1200+175, 650+100, 760+400, 300+100),# Cloud → Central
    (300+175, 650+100, 760, 300+100)      # Election → Central
]

for x1, y1, x2, y2 in connections:
    line = slide.shapes.add_connector(
        1,  # straight line
        px(x1), px(y1),
        px(x2), px(y2)
    )
    line.line.color.rgb = RGBColor(80, 80, 80)
    line.line.width = Pt(2)

# --- Footer ---
footer = slide.shapes.add_textbox(px(400), px(900), px(1200), px(60))
tf = footer.text_frame
tf.text = "A multi-actor ecosystem shaping national digital resilience and sovereignty."
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(85, 85, 85)
p.font.name = "Segoe UI"

prs.save("ecosystem_map_generated.pptx")
